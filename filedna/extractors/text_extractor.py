"""
FileDNA – text extraction engine.

extract_text(path) → str

Pulls clean, readable plain text from any supported file type.
This is what every RAG pipeline needs before chunking + embedding.

Supported:
  Documents : pdf, docx, pptx, xlsx, csv, txt, md, html, xml, json, epub
  Images    : png, jpg, gif, bmp, tiff, webp  (returns "" — no OCR without AI)
  Audio     : mp3, wav, m4a, flac, ogg        (returns "" — needs AI transcription)
  Video     : mp4, mov, mkv, webm, avi        (returns "" — needs AI transcription)
  Archives  : zip, tar, gz                    (extracts + recurses into text files)

Design:
  - Each extractor is a standalone function — easy to test and replace
  - Never raises — returns "" on failure (caller decides what to do)
  - All imports are lazy — only loaded for the relevant type
  - ExtractionResult carries the text AND metadata about how it was extracted
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any


@dataclass
class ExtractionResult:
    """Result of a text extraction operation."""
    text: str = ""
    char_count: int = 0
    word_count: int = 0
    extractor: str = "unknown"          # which extractor was used
    truncated: bool = False             # True if text was cut due to max_chars
    warnings: list[str] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)

    def __post_init__(self) -> None:
        if self.text:
            self.char_count = len(self.text)
            self.word_count = len(self.text.split())

    def __bool__(self) -> bool:
        return bool(self.text.strip())


# ---------------------------------------------------------------------------
# Document extractors
# ---------------------------------------------------------------------------

def _extract_pdf(path: Path, max_chars: int) -> ExtractionResult:
    try:
        import pdfplumber  # type: ignore
        parts: list[str] = []
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    parts.append(t)
        text = "\n\n".join(parts)
        truncated = False
        if max_chars and len(text) > max_chars:
            text = text[:max_chars]
            truncated = True
        return ExtractionResult(text=text, extractor="pdfplumber", truncated=truncated)
    except Exception as exc:
        return ExtractionResult(warnings=[f"PDF extraction failed: {exc}"])


def _extract_docx(path: Path, max_chars: int) -> ExtractionResult:
    try:
        from docx import Document  # type: ignore
        doc = Document(str(path))
        parts: list[str] = []
        # Body paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    parts.append(row_text)
        text = "\n".join(parts)
        truncated = False
        if max_chars and len(text) > max_chars:
            text = text[:max_chars]
            truncated = True
        return ExtractionResult(text=text, extractor="python-docx", truncated=truncated)
    except Exception as exc:
        return ExtractionResult(warnings=[f"DOCX extraction failed: {exc}"])


def _extract_pptx(path: Path, max_chars: int) -> ExtractionResult:
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(path))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_parts.append(t)
            if slide_parts:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_parts))
        text = "\n\n".join(parts)
        truncated = False
        if max_chars and len(text) > max_chars:
            text = text[:max_chars]
            truncated = True
        return ExtractionResult(text=text, extractor="python-pptx", truncated=truncated)
    except Exception as exc:
        return ExtractionResult(warnings=[f"PPTX extraction failed: {exc}"])


def _extract_xlsx(path: Path, max_chars: int) -> ExtractionResult:
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts: list[str] = []
        for ws in wb.worksheets:
            sheet_rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    sheet_rows.append(" | ".join(cells))
            if sheet_rows:
                parts.append(f"[Sheet: {ws.title}]\n" + "\n".join(sheet_rows))
        wb.close()
        text = "\n\n".join(parts)
        truncated = False
        if max_chars and len(text) > max_chars:
            text = text[:max_chars]
            truncated = True
        return ExtractionResult(text=text, extractor="openpyxl", truncated=truncated)
    except Exception as exc:
        return ExtractionResult(warnings=[f"XLSX extraction failed: {exc}"])


def _extract_csv(path: Path, max_chars: int) -> ExtractionResult:
    try:
        import csv
        rows: list[str] = []
        with open(path, newline="", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" | ".join(row))
        text = "\n".join(rows)
        truncated = False
        if max_chars and len(text) > max_chars:
            text = text[:max_chars]
            truncated = True
        return ExtractionResult(text=text, extractor="csv", truncated=truncated)
    except Exception as exc:
        return ExtractionResult(warnings=[f"CSV extraction failed: {exc}"])


def _extract_html(path: Path, max_chars: int) -> ExtractionResult:
    try:
        # trafilatura gives much cleaner output than BeautifulSoup for HTML
        import trafilatura  # type: ignore
        raw = path.read_text(encoding="utf-8", errors="replace")
        text = trafilatura.extract(raw) or ""
        if not text:
            # fallback: strip tags manually
            import re
            text = re.sub(r"<[^>]+>", " ", raw)
            text = re.sub(r"\s+", " ", text).strip()
        truncated = False
        if max_chars and len(text) > max_chars:
            text = text[:max_chars]
            truncated = True
        return ExtractionResult(text=text, extractor="trafilatura", truncated=truncated)
    except Exception as exc:
        return ExtractionResult(warnings=[f"HTML extraction failed: {exc}"])


def _extract_xml(path: Path, max_chars: int) -> ExtractionResult:
    try:
        import xml.etree.ElementTree as ET
        tree = ET.parse(str(path))
        parts: list[str] = []
        for elem in tree.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                parts.append(elem.tail.strip())
        text = "\n".join(parts)
        truncated = False
        if max_chars and len(text) > max_chars:
            text = text[:max_chars]
            truncated = True
        return ExtractionResult(text=text, extractor="xml.etree", truncated=truncated)
    except Exception as exc:
        return ExtractionResult(warnings=[f"XML extraction failed: {exc}"])


def _extract_json(path: Path, max_chars: int) -> ExtractionResult:
    try:
        import json
        with open(path, encoding="utf-8", errors="replace") as f:
            data = json.load(f)
        # Pretty-print so nested keys are readable
        text = json.dumps(data, indent=2, ensure_ascii=False)
        truncated = False
        if max_chars and len(text) > max_chars:
            text = text[:max_chars]
            truncated = True
        return ExtractionResult(text=text, extractor="json", truncated=truncated)
    except Exception as exc:
        return ExtractionResult(warnings=[f"JSON extraction failed: {exc}"])


def _extract_epub(path: Path, max_chars: int) -> ExtractionResult:
    try:
        import zipfile
        import re
        parts: list[str] = []
        with zipfile.ZipFile(path) as z:
            for name in sorted(z.namelist()):
                if name.endswith((".html", ".xhtml", ".htm")):
                    raw = z.read(name).decode("utf-8", errors="replace")
                    text = re.sub(r"<[^>]+>", " ", raw)
                    text = re.sub(r"\s+", " ", text).strip()
                    if text:
                        parts.append(text)
        text = "\n\n".join(parts)
        truncated = False
        if max_chars and len(text) > max_chars:
            text = text[:max_chars]
            truncated = True
        return ExtractionResult(text=text, extractor="epub+zip", truncated=truncated)
    except Exception as exc:
        return ExtractionResult(warnings=[f"EPUB extraction failed: {exc}"])


def _extract_plain(path: Path, max_chars: int) -> ExtractionResult:
    """For .txt, .md, and any other plain text type."""
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        truncated = False
        if max_chars and len(text) > max_chars:
            text = text[:max_chars]
            truncated = True
        return ExtractionResult(text=text, extractor="plain-text", truncated=truncated)
    except Exception as exc:
        return ExtractionResult(warnings=[f"Plain text extraction failed: {exc}"])


def _extract_archive(path: Path, max_chars: int) -> ExtractionResult:
    """
    Recurse into ZIP/TAR/GZ archives and extract text from any text files inside.
    Only goes one level deep — no nested archives.
    """
    import zipfile, tarfile, gzip, io

    TEXT_EXTENSIONS = {
        ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm",
        ".py", ".js", ".ts", ".css", ".yaml", ".yml", ".toml", ".ini", ".cfg",
        ".rst", ".log", ".sql",
    }

    parts: list[str] = []
    warnings: list[str] = []
    suffix = path.suffix.lower()

    try:
        if suffix in (".zip",):
            with zipfile.ZipFile(path) as z:
                for info in z.infolist():
                    if Path(info.filename).suffix.lower() in TEXT_EXTENSIONS:
                        try:
                            content = z.read(info.filename).decode("utf-8", errors="replace")
                            parts.append(f"[{info.filename}]\n{content}")
                        except Exception:
                            pass
        elif suffix in (".tar",):
            with tarfile.open(str(path)) as t:
                for member in t.getmembers():
                    if member.isfile() and Path(member.name).suffix.lower() in TEXT_EXTENSIONS:
                        try:
                            f = t.extractfile(member)
                            if f:
                                content = f.read().decode("utf-8", errors="replace")
                                parts.append(f"[{member.name}]\n{content}")
                        except Exception:
                            pass
        elif suffix in (".gz",):
            try:
                with gzip.open(str(path), "rb") as f:
                    content = f.read().decode("utf-8", errors="replace")
                    parts.append(content)
            except Exception:
                pass
    except Exception as exc:
        warnings.append(f"Archive extraction failed: {exc}")

    text = "\n\n".join(parts)
    truncated = False
    if max_chars and len(text) > max_chars:
        text = text[:max_chars]
        truncated = True

    if not text and not warnings:
        warnings.append("No readable text files found inside archive")

    return ExtractionResult(
        text=text,
        extractor="archive",
        truncated=truncated,
        warnings=warnings,
    )


# ---------------------------------------------------------------------------
# Dispatch table
# ---------------------------------------------------------------------------

_DISPATCH = {
    "pdf":  _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "xlsx": _extract_xlsx,
    "csv":  _extract_csv,
    "html": _extract_html,
    "xml":  _extract_xml,
    "json": _extract_json,
    "epub": _extract_epub,
    "txt":  _extract_plain,
    "md":   _extract_plain,
    "zip":  _extract_archive,
    "tar":  _extract_archive,
    "gz":   _extract_archive,
}

# These types produce no text without AI assistance — return empty with explanation
_NO_TEXT_TYPES = {
    "png", "jpg", "jpeg", "gif", "bmp", "tiff", "tif", "webp", "svg",
    "mp3", "wav", "flac", "ogg", "m4a", "aac",
    "mp4", "mov", "mkv", "webm", "avi",
}


def extract_text(path: Path, real_type: str, max_chars: int = 0) -> ExtractionResult:
    """
    Extract plain text from a file.

    Args:
        path:      Path to the file.
        real_type: Already-detected real type string (e.g. "pdf").
        max_chars: Truncate text at this many characters (0 = no limit).

    Returns:
        ExtractionResult — always succeeds, never raises.
    """
    if real_type in _NO_TEXT_TYPES:
        return ExtractionResult(
            warnings=[
                f"Type '{real_type}' requires AI assistance for text extraction. "
                "Use filedna.ai.transcribe() for audio/video, or filedna.ai.ocr() for images."
            ]
        )

    fn = _DISPATCH.get(real_type)
    if fn is None:
        # Try plain text as a last resort
        try:
            return _extract_plain(path, max_chars)
        except Exception:
            return ExtractionResult(warnings=[f"No extractor for type '{real_type}'"])

    try:
        return fn(path, max_chars)
    except Exception as exc:
        return ExtractionResult(warnings=[f"Extraction error: {exc}"])
