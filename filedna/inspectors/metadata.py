"""
FileDNA – metadata inspectors.

Each inspector takes a Path and returns a dict of metadata.
Inspectors are lazy (only imported when needed) and never raise –
they return whatever they can extract.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

MetaDict = dict[str, Any]

# ---------------------------------------------------------------------------
# Token estimation helper — encoder cached at module level (not per-call)
# ---------------------------------------------------------------------------

_tiktoken_encoder = None

def _get_encoder():
    global _tiktoken_encoder
    if _tiktoken_encoder is None:
        try:
            import tiktoken  # type: ignore
            _tiktoken_encoder = tiktoken.get_encoding("cl100k_base")
        except Exception:
            pass
    return _tiktoken_encoder


def _estimate_tokens(text: str) -> int:
    """Estimate token count. Uses cached tiktoken encoder, else word heuristic."""
    enc = _get_encoder()
    if enc is not None:
        try:
            return len(enc.encode(text))
        except Exception:
            pass
    # Fallback: ~0.75 tokens per word
    return max(1, int(len(text.split()) / 0.75))


# ---------------------------------------------------------------------------
# Language detection — uses langdetect with fixed seed (deterministic, 3MB)
# ---------------------------------------------------------------------------

def _detect_language(text: str) -> str:
    """Detect language of text. Returns ISO 639-1 code or 'unknown'."""
    if not text.strip():
        return "unknown"
    try:
        from langdetect import detect, DetectorFactory  # type: ignore
        DetectorFactory.seed = 0   # makes results deterministic (same text = same result)
        return detect(text[:5000])
    except Exception:
        return "unknown"


# ---------------------------------------------------------------------------
# Human-readable size
# ---------------------------------------------------------------------------

def human_size(n: int) -> str:
    for unit in ("B", "KB", "MB", "GB", "TB"):
        if n < 1024:
            return f"{n:.1f} {unit}" if unit != "B" else f"{n} {unit}"
        n /= 1024  # type: ignore[assignment]
    return f"{n:.1f} PB"


# ---------------------------------------------------------------------------
# Document inspectors
# ---------------------------------------------------------------------------

def inspect_pdf(path: Path) -> MetaDict:
    meta: MetaDict = {}
    try:
        import pdfplumber  # type: ignore
        with pdfplumber.open(str(path)) as pdf:
            meta["pages"] = len(pdf.pages)
            meta["encrypted"] = False  # pdfplumber already handles decryption

            all_text: list[str] = []
            has_images = False
            has_tables = False

            for page in pdf.pages:
                text = page.extract_text() or ""
                all_text.append(text)
                if page.images:
                    has_images = True
                if page.extract_tables():
                    has_tables = True

            full_text = "\n".join(all_text)
            meta["contains_images"] = has_images
            meta["contains_tables"] = has_tables
            meta["language"] = _detect_language(full_text)
            meta["estimated_tokens"] = _estimate_tokens(full_text)
    except Exception as exc:
        meta["inspection_error"] = str(exc)
    return meta


def inspect_docx(path: Path) -> MetaDict:
    meta: MetaDict = {}
    try:
        from docx import Document  # type: ignore
        doc = Document(str(path))
        paragraphs = len(doc.paragraphs)
        words = sum(len(p.text.split()) for p in doc.paragraphs)
        full_text = "\n".join(p.text for p in doc.paragraphs)
        meta["paragraphs"] = paragraphs
        meta["words"] = words
        # Rough page estimate: ~300 words per page
        meta["estimated_pages"] = max(1, round(words / 300))
        meta["language"] = _detect_language(full_text)
        meta["estimated_tokens"] = _estimate_tokens(full_text)
    except Exception as exc:
        meta["inspection_error"] = str(exc)
    return meta


def inspect_xlsx(path: Path) -> MetaDict:
    meta: MetaDict = {}
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        sheets = wb.sheetnames
        meta["sheets"] = len(sheets)
        meta["sheet_names"] = sheets

        total_rows = 0
        total_cols = 0
        all_text: list[str] = []

        for ws in wb.worksheets:
            total_rows += ws.max_row or 0
            total_cols = max(total_cols, ws.max_column or 0)
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        all_text.append(str(cell))

        wb.close()
        meta["rows"] = total_rows
        meta["columns"] = total_cols
        meta["estimated_tokens"] = _estimate_tokens(" ".join(all_text))
    except Exception as exc:
        meta["inspection_error"] = str(exc)
    return meta


def inspect_pptx(path: Path) -> MetaDict:
    meta: MetaDict = {}
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(path))
        slides = len(prs.slides)
        all_text: list[str] = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            all_text.append(t)

        full_text = "\n".join(all_text)
        meta["slides"] = slides
        meta["text_length"] = len(full_text)
        meta["language"] = _detect_language(full_text)
        meta["estimated_tokens"] = _estimate_tokens(full_text)
    except Exception as exc:
        meta["inspection_error"] = str(exc)
    return meta


def inspect_epub(path: Path) -> MetaDict:
    meta: MetaDict = {}
    try:
        import zipfile
        with zipfile.ZipFile(path) as z:
            text_parts: list[str] = []
            for name in z.namelist():
                if name.endswith((".html", ".xhtml", ".htm")):
                    raw = z.read(name).decode("utf-8", errors="replace")
                    # Strip tags naively
                    import re
                    text_parts.append(re.sub(r"<[^>]+>", " ", raw))
            full_text = " ".join(text_parts)
            meta["language"] = _detect_language(full_text)
            meta["estimated_tokens"] = _estimate_tokens(full_text)
    except Exception as exc:
        meta["inspection_error"] = str(exc)
    return meta


def _inspect_text_file(path: Path) -> MetaDict:
    meta: MetaDict = {}
    try:
        with open(path, encoding="utf-8", errors="replace") as f:
            text = f.read()
        meta["characters"] = len(text)
        meta["lines"] = text.count("\n") + 1
        meta["words"] = len(text.split())
        meta["language"] = _detect_language(text)
        meta["estimated_tokens"] = _estimate_tokens(text)
    except Exception as exc:
        meta["inspection_error"] = str(exc)
    return meta


def inspect_csv(path: Path) -> MetaDict:
    meta: MetaDict = {}
    try:
        import csv
        with open(path, newline="", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            rows = list(reader)
        meta["rows"] = len(rows)
        meta["columns"] = max((len(r) for r in rows), default=0)
        flat = " ".join(" ".join(r) for r in rows)
        meta["estimated_tokens"] = _estimate_tokens(flat)
    except Exception as exc:
        meta["inspection_error"] = str(exc)
    return meta


# ---------------------------------------------------------------------------
# Image inspectors
# ---------------------------------------------------------------------------

def inspect_image(path: Path) -> MetaDict:
    meta: MetaDict = {}
    try:
        from PIL import Image  # type: ignore
        with Image.open(str(path)) as img:
            meta["width"] = img.width
            meta["height"] = img.height
            meta["mode"] = img.mode
            dpi = img.info.get("dpi")
            if dpi:
                meta["dpi"] = (round(dpi[0]), round(dpi[1]))
            meta["has_transparency"] = img.mode in ("RGBA", "LA", "PA") or (
                "transparency" in img.info
            )
    except Exception as exc:
        meta["inspection_error"] = str(exc)
    return meta


# ---------------------------------------------------------------------------
# Audio inspectors
# ---------------------------------------------------------------------------

def inspect_audio(path: Path) -> MetaDict:
    meta: MetaDict = {}
    try:
        import mutagen  # type: ignore
        tag = mutagen.File(str(path))
        if tag is not None:
            info = tag.info
            meta["duration"] = round(getattr(info, "length", 0), 2)
            meta["bitrate"] = getattr(info, "bitrate", None)
            meta["sample_rate"] = getattr(info, "sample_rate", None)
            meta["channels"] = getattr(info, "channels", None)
    except Exception as exc:
        meta["inspection_error"] = str(exc)
    return meta


# ---------------------------------------------------------------------------
# Video inspectors
# ---------------------------------------------------------------------------

def inspect_video(path: Path) -> MetaDict:
    meta: MetaDict = {}

    # Try ffprobe first
    try:
        import json
        import subprocess
        result = subprocess.run(
            [
                "ffprobe", "-v", "quiet",
                "-print_format", "json",
                "-show_streams", "-show_format",
                str(path),
            ],
            capture_output=True, text=True, timeout=10
        )
        if result.returncode == 0:
            data = json.loads(result.stdout)
            fmt = data.get("format", {})
            meta["duration"] = round(float(fmt.get("duration", 0)), 2)
            for stream in data.get("streams", []):
                if stream.get("codec_type") == "video":
                    meta["width"] = stream.get("width")
                    meta["height"] = stream.get("height")
                    meta["resolution"] = f"{stream.get('width')}x{stream.get('height')}"
                    meta["codec"] = stream.get("codec_name")
                    r_frame = stream.get("avg_frame_rate", "0/1")
                    try:
                        num, den = r_frame.split("/")
                        meta["fps"] = round(int(num) / int(den), 2)
                    except Exception:
                        pass
                    break
            return meta
    except (FileNotFoundError, Exception):
        pass

    # Fallback: mutagen for duration on mp4/m4a
    try:
        import mutagen  # type: ignore
        tag = mutagen.File(str(path))
        if tag is not None and hasattr(tag, "info"):
            meta["duration"] = round(getattr(tag.info, "length", 0), 2)
    except Exception:
        pass

    if not meta:
        meta["inspection_note"] = "Install ffprobe for full video metadata"

    return meta


# ---------------------------------------------------------------------------
# Archive inspectors
# ---------------------------------------------------------------------------

def inspect_zip(path: Path) -> MetaDict:
    meta: MetaDict = {}
    import zipfile
    try:
        with zipfile.ZipFile(path) as z:
            members = z.infolist()
            meta["file_count"] = len(members)
            meta["total_uncompressed_bytes"] = sum(m.file_size for m in members)
            meta["total_compressed_bytes"] = sum(m.compress_size for m in members)
    except Exception as exc:
        meta["inspection_error"] = str(exc)
    return meta


def inspect_tar(path: Path) -> MetaDict:
    meta: MetaDict = {}
    import tarfile
    try:
        with tarfile.open(str(path)) as t:
            members = t.getmembers()
            meta["file_count"] = len(members)
            meta["total_bytes"] = sum(m.size for m in members if m.isfile())
    except Exception as exc:
        meta["inspection_error"] = str(exc)
    return meta


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

_IMAGE_TYPES = {"png", "jpg", "jpeg", "gif", "bmp", "webp", "tiff", "tif", "svg"}
_AUDIO_TYPES = {"mp3", "wav", "flac", "ogg", "m4a", "aac"}
_VIDEO_TYPES = {"mp4", "mov", "mkv", "webm", "avi"}
_TEXT_TYPES  = {"txt", "md", "html", "xml", "json"}

_DISPATCH = {
    "pdf":  inspect_pdf,
    "docx": inspect_docx,
    "xlsx": inspect_xlsx,
    "pptx": inspect_pptx,
    "epub": inspect_epub,
    "csv":  inspect_csv,
    "zip":  inspect_zip,
    "tar":  inspect_tar,
}
for _t in _IMAGE_TYPES:
    _DISPATCH[_t] = inspect_image
for _t in _AUDIO_TYPES:
    _DISPATCH[_t] = inspect_audio
for _t in _VIDEO_TYPES:
    _DISPATCH[_t] = inspect_video
for _t in _TEXT_TYPES:
    _DISPATCH[_t] = _inspect_text_file


def inspect(path: Path, real_type: str) -> MetaDict:
    """Extract metadata from *path* for *real_type*. Never raises."""
    fn = _DISPATCH.get(real_type)
    if fn is None:
        return {}
    try:
        return fn(path)
    except Exception as exc:
        return {"inspection_error": str(exc)}